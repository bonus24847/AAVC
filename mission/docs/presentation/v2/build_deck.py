# -*- coding: utf-8 -*-
"""One content spec -> HTML deck + PPTX deck (same 1280x720 canvas)."""
import base64, html as H, os, re, sys, shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

S=os.environ['SCRATCH']; A=S+'/assets'; OUT=os.environ['OUT']
os.makedirs(OUT+'/media',exist_ok=True)
C=dict(bg='0a0e15',bg2='0e1420',surf='141d2b',line='26344a',line2='33455f',ink='e9eff7',mut='93a4bd',dim='65788f',acc='33e3c4',amb='f0b429',go='3fce7c',crit='ff5f66')
def rgb(k): return RGBColor.from_string(C[k])

# ---------------- block constructors ----------------
def img(src,x,y,w,h,cap=None): return dict(t='img',src=src,x=x,y=y,w=w,h=h,cap=cap)
def ph(text,x,y,w,h,sub=None): return dict(t='ph',text=text,x=x,y=y,w=w,h=h,sub=sub)
def txt(lines,x,y,w,h,size=17): return dict(t='txt',lines=lines,x=x,y=y,w=w,h=h,size=size)
def tiles(items,x,y,w,h,cols=None): return dict(t='tiles',items=items,x=x,y=y,w=w,h=h,cols=cols or len(items))
def video(src,poster,x,y,w,h,cap=None): return dict(t='video',src=src,poster=poster,x=x,y=y,w=w,h=h,cap=cap)
def table(rows,x,y,w,h,colw,size=13): return dict(t='table',rows=rows,x=x,y=y,w=w,h=h,colw=colw,size=size)
def box(x,y,w,h,color='line'): return dict(t='box',x=x,y=y,w=w,h=h,color=color)

# ---------------- SLIDES ----------------
SL=[]
def slide(eyebrow,title,blocks,notes='',kind='normal'): SL.append(dict(eyebrow=eyebrow,title=title,blocks=blocks,notes=notes,kind=kind))

# 1 TITLE
slide('AAVC 2026 · IAAI KMITL · 28–30 AUG 2026','',[
  txt(['## AeroOptix','**King Mongkut\'s University of Technology North Bangkok** · Faculty of Engineering',
       '',
       '**Sky-Field Egg Delivery** — a fully autonomous PX4 hexacopter that finds the committee-assigned ArUco pads and lands to deliver fragile cargo, hands-off, inside the 20-minute window.',
       '',
       '**FIRST-TIME ENTRANT** — our first AAVC, built from zero in four months.'],60,150,620,300,17),
  txt(['## Team','• [ชื่อ] — team leader / mission software','• [ชื่อ] — flight systems / safety pilot','• [ชื่อ] — vision & GCS','• [ชื่อ] — airframe & power','• Advisor: [อาจารย์ที่ปรึกษา]'],60,470,620,220,15),
  ph('รูปทีม + โดรน (รอถ่าย)',720,150,500,470,'team photo with the aircraft'),
],'สวัสดีครับ ทีม AeroOptix จาก มจพ. — นี่คือการแข่ง AAVC ครั้งแรกของพวกเรา เริ่มจากศูนย์เมื่อ 4 เดือนก่อน แนะนำสมาชิกและบทบาท','title')

# 2 MISSION RECAP
slide('01 · MISSION RECAP (RULES V1.3)','What the organiser asked for',[
  img(A+'/rules_map_transit.jpg',60,130,640,388,'Rules Fig. 3 — controlled airspace (red), search area (yellow), no-fly (orange), L&R + transit route'),
  img(A+'/rules_pad.jpg',60,530,200,160,'pad 1×1 m · ring ⌀750 · ArUco 400 mm'),
  txt(['• **6 landing pads** in the search area, each a 1×1 m white square with a **400 mm ArUco marker** (DICT_4X4_50, ids 1–6)',
       '• The committee **assigns 4 of the 6** to our team at GO — the other 2 are decoys',
       '• Every flight must fly the **mandatory transit P1→P2→P3 at 20 m**, both ways, scored per point',
       '• **Land ON the pad**, then release that pad\'s egg (heart box) — never airborne',
       '• Everything inside **one 20-minute operation window**; ceiling 20 m, below 10 m only for the delivery descent',
       '• **Full autonomy = full marks** — any manual intervention costs points'],280,530,940,160,13),
  tiles([('6→4','pads · assigned','2 decoys ignored','acc'),('20 m','transit P1·P2·P3','strictly, both ways','ink'),('20:00','operation window','no bonus for early','amb'),('10–20 m','search band','descend only over pad','ink'),('AUTO','top marks','no intervention','go')],720,130,500,388,1),
],'สรุปโจทย์จากผู้จัด: แผนที่สนาม KMITL — พื้นที่ควบคุม, search area, no-fly, จุด L&R และ transit P1-P3 ที่บังคับบินที่ 20 m; pad 6 แผ่น assign 4; ลงจอดบนแผ่นแล้วปล่อยไข่; 20 นาที; autonomy เต็ม = คะแนนเต็ม')

# 3 FLIGHT PATH
slide('02 · OUR FLIGHT PATH','Survey everything once · deliver each egg · return home',[
  img(A+'/profile.png',60,128,1160,262,None),
  img(A+'/kmitl_plan.png',60,400,700,290,'Planned flight at KMITL — built from the rules coordinates by our sweep planner'),
  txt(['## Why this shape — time comes first',
       '• A heart must arrive **fast and safe** → mission time is our first design driver',
       '• **Survey the whole area once** (5-leg sweep at 12 m), decoding every pad into a registry',
       '• Then **serve the 4 assigned pads by the shortest route** — land, release, climb, next',
       '• Egress through the corridor, land at L&R, disarm',
       '• Budget: **~900 s of 1200** for a full 4-egg flight (SITL: 536 s)'],780,400,440,290,14),
],'เส้นทางบิน: takeoff → P1 P2 P3 ที่ 20 m → sweep 12 m ให้ครบทั้งพื้นที่ในรอบเดียว → เสิร์ฟ 4 pad ตามเส้นทางสั้นสุด → กลับผ่าน corridor → ลง L&R. เหตุผล: ส่งหัวใจต้องเร็วและปลอดภัย เวลาจึงเป็นตัวแปรแรกที่พิจารณา — กวาดครั้งเดียวดีกว่าหาทีละ pad')

# 4 AIRCRAFT
slide('03 · THE AIRCRAFT','EFT X6100 hexacopter — Pixhawk 6X + Raspberry Pi CM4',[
  ph('รูปโดรน EFT X6100 (รอถ่าย)',60,128,520,380,'aircraft on the pad, eggs loaded'),
  table([['Subsystem','What we fly'],
         ['Airframe','EFT X6100 hexacopter, 18″ props, carbon — 1.0 m wheelbase'],
         ['Flight controller','Holybro Pixhawk 6X · PX4 v1.17'],
         ['Companion','Raspberry Pi CM4 — mission brain, vision, GCS link'],
         ['Camera','OV9281 mono global-shutter, nadir, 74.2° lens (measured)'],
         ['Height','Benewake TFmini-S lidar + barometer (EKF2 baro ref)'],
         ['Radio','RadioMaster NOMAD — RC + MAVLink telemetry in one link'],
         ['Power','2 × 6S packs (17 000 + 15 000 mAh) · PM02D → FC'],
         ['Payload','4 egg latches on AUX 4/1/2/3, released in diagonal order']],600,128,620,380,[150,470],12),
  tiles([('7.2 kg','all-up weight','29 % of 25 kg MTOW','acc'),('~31 A','hover current','6 rotors, inside rating','ink'),('5 m/s','max cruise · 3 m/s sweep','MPC_XY_VEL_MAX','ink'),('~4.5 Ah','per 4-egg flight','of 12.75 Ah usable','amb')],60,520,640,170,4),
  txt(['## Redundancy by design',
       '• **Six rotors** — with one rotor out PX4\'s allocator still holds attitude and thrust; only yaw is given up, and the safety pilot recovers it',
       '• **Two battery packs** — a full spare pack for the second flight / a swap between flights',
       '• **Height from three sources** — lidar, barometer, GPS'],720,520,500,170,13),
],'โดรน EFT X6100 hexacopter (ได้รับความอนุเคราะห์จาก DRON) + Pixhawk 6X + CM4 อุปกรณ์ตามตาราง สมรรถนะ: 7.2 kg, hover ~31 A, ใช้ ~4.5 Ah ต่อเที่ยว 4 ไข่ ข้อดี 6 ใบพัด = redundant ถ้าใบพัดเสีย 1 ตัวยังบินได้; แบต 2 ก้อนเป็น redundant')

# 5.1 ARCHITECTURE
slide('04 · SYSTEM ARCHITECTURE','PX4 firmware · MAVLink protocol · one deterministic loop',[
  img(A+'/arch.png',60,128,780,234,None),
  txt(['## Firmware & protocol',
       '• **PX4 v1.17** on the Pixhawk 6X flies the aircraft and owns every failsafe',
       '• **MAVLink** everywhere: MAVSDK on the CM4 → FC over serial; pymavlink on the GCS over the NOMAD radio',
       '• **CM4 mission brain** (Python, async) — plan, search, vision, landing, audit — **offline, no cloud, no LLM**',
       '• **OpenCV ArUco** decode of the pad markers — classical CV, deterministic',
       '• **Safety pilot RC** (ELRS) can retake the aircraft at any instant'],860,128,360,420,13),
  tiles([('PX4 1.17','firmware','Pixhawk 6X','acc'),('MAVLink','protocol','MAVSDK · pymavlink','ink'),('CM4','mission brain','Python · OpenCV','amb'),('NOMAD','RC + telemetry','one radio link','ink')],60,380,780,150,4),
  txt(['**Air segment:** Pixhawk 6X · nadir camera · TFmini-S · NOMAD radio · egg rack ×4 · 6 × PWM ESC · CM4.  **Ground segment:** RC (safety pilot) · our web GCS (telemetry · imaging · status) — the GCS monitors, it never has to command.'],60,545,780,140,13),
],'โครงสร้างระบบ: air segment มี Pixhawk 6X (PX4 1.17) เป็น flight controller, CM4 เป็นสมองภารกิจ, เชื่อมกันด้วย MAVLink (MAVSDK); ground segment มี RC ของ safety pilot กับ GCS ที่เราพัฒนาเองผ่านวิทยุ NOMAD (pymavlink)')

# 5.2 FLOWCHART
slide('05 · FULL AUTONOMY','Two human actions, then the CM4 flies the whole mission',[
  img(S+'/diag/flow.png',60,128,1160,562,None),
],'flowchart การบินอัตโนมัติ: มนุษย์ทำแค่ 2 อย่าง (stage id + arm/OFFBOARD) หลังจากนั้น CM4 บินเอง: takeoff → transit → sweep → registry → order → loop ต่อ pad (gate เวลา/แบต, id verify, land, touchdown, release) → transit กลับ → land → disarm; watchdog ทำงานตลอด')

# 6a SITL video
slide('06 · SOFTWARE-IN-THE-LOOP','Our own SITL — the mission flew in simulation first',[
  video('media/sitl_replay.mp4',A+'/sitl_poster.jpg',60,128,640,360,'SITL replay ×5 — nadir camera, phase, AGL, battery, track inset (PX4 v1.17 + Gazebo, 6 ArUco pads, ids randomised)'),
  img(A+'/sitl_track.png',720,128,500,330,'Flown track — transit, 5-leg sweep, 4 land-on-pad deliveries, egress'),
  txt(['• **PX4 v1.17 + Gazebo** with a model of the KMUTNB practice field, our EFT X6100 model and 6 pad models — ids and positions re-rolled per run',
       '• The **same orchestrator code** that runs on the CM4 flies the simulator — no sim-only path',
       '• Flown on the **competition envelope** (transit 20 m, sweep 12 m → 28 px marker): **4/4 eggs, 6/6 pads decoded, 0.09–0.22 m, 536 s**',
       '• Every run is checked by an independent post-flight verifier that **fails closed**'],60,505,1160,185,14),
],'SITL ที่เราพัฒนาเอง: PX4 + Gazebo + โมเดลสนามซ้อม + pad 6 แผ่น โค้ดภารกิจตัวเดียวกับที่รันบน CM4 วิดีโอคือ replay ×5 จากกล้อง nadir ใน sim พร้อม phase/AGL/แบต; แผนที่ขวาคือเส้นทางที่บินจริงใน sim')

# 6b SITL charts
slide('07 · SITL ANALYSIS','What the simulator taught us before the first real flight',[
  img(A+'/sitl_altitude.png',60,128,720,300,None),
  txt(['## Reading the graphs',
       '• Transit holds **19.5 m** under the 20 m ceiling; the sweep sits at **12 m**, above the 10 m floor',
       '• Four descents to **0 m** — each release only after touchdown',
       '• Battery never approaches the **30 % return floor**; the energy gate stays silent',
       '• Releases land **0.10–0.16 m** from pad centre on a 1 m pad',
       '• Sweep spacing was **A/B-tested**: 0.30 overlap keeps 6/6 coverage and saves 111 s'],800,128,420,300,13),
  img(A+'/sitl_battery.png',60,440,570,250,None),
  img(A+'/sitl_results.png',650,440,570,250,None),
],'กราฟจาก SITL: ความสูงตามเฟส (transit 19.5 m, sweep 12 m, ลง 4 ครั้ง), แบตไม่แตะ floor 30%, ความแม่นปล่อย 0.10-0.16 m, และการทดสอบ A/B ระยะห่างแนวกวาด 0.30 vs 0.44 — coverage เท่ากัน เร็วกว่า 111 วินาที')

# 7 GCS
slide('08 · GROUND CONTROL STATION','Built by us for this mission — a web GCS on the radio link',[
  img(A+'/gcs_demo.jpg',60,128,820,461,'AAVC GCS (demo mode) — live map with rulebook zones, mission panel, sensor & link chips, 6-motor bars'),
  txt(['• **Python web GCS** (http.server + pymavlink) — no QGroundControl in the loop',
       '• **Map** with the controlled airspace, search area and transit; the drone and every **decoded pad** plotted (assigned / mapped / delivered)',
       '• **Pad assignment** — tick the 4 committee ids, shown as the **ArUco glyph** to match the card picture-to-picture',
       '• **Mission panel** — phase progress, pad status, the **20:00 window clock**',
       '• Sensor / RC / link / camera / CM4 health chips, 6-motor bars, attitude & battery',
       '• Draws only from the **NOMAD radio** — a WiFi dropout never freezes a panel'],900,128,320,560,12),
  ph('ภาพ GCS ขณะบินจริง (รอ capture)',60,600,820,90,'field screenshot'),
],'GCS ที่เราพัฒนาขึ้นเอง: web GCS ภาษา Python ผ่าน pymavlink — แผนที่พร้อมโซนจากกติกา, pad ที่ decode ได้, เลือก id ที่ได้รับ assign เป็นรูป marker, mission panel กับนาฬิกา 20 นาที, chip สถานะเซนเซอร์/ลิงก์, ทำงานผ่านวิทยุอย่างเดียว')

# 8 SAFETY
slide('09 · FAILSAFE · CHECKLIST · REDUNDANCY','Three independent safety layers — the pilot always wins',[
  txt(['## 1 · Safety pilot (RC)','Retakes the aircraft on POSCTL at any instant; the companion **stands down** within 1 s and refuses every command — drilled on the real FC (21 Aug).',
       '## 2 · Flight controller (PX4)','Geofence → **Return** (value verified on the board) · datalink loss → Return · RC loss → Return · low battery → Return then Land · 50 m vertical net.',
       '## 3 · Companion watchdog','Ceiling 20 m / floor 10 m on AGL latched at arming · no-fly · battery floor · GPS · telemetry age · detects an FC-initiated RTL/LAND and stands down.'],60,128,420,560,12),
  img(A+'/checklist.jpg',500,128,300,424,'Competition checklist — IAAI KMITL (p.1 of 6)'),
  txt(['## Before every flight','• **24 board parameters read back** — motor map, battery gauge, fence, height source (a motor map once vanished; we check, not trust)',
       '• GPS matched to FC · inside the fence · battery · link · RC · egg latches',
       '• Formal checklist: schedule · standby silence · safety inspection · board gate · GO two-person rule · emergency · post-flight'],500,560,300,130,11),
  txt(['## Redundancy','• **6 rotors** — one out: attitude & thrust held, pilot takes yaw',
       '• **2 battery packs** — spare / swap between flights; gauge reads conservative',
       '• **Height ×3** — TFmini lidar, barometer (EKF reference), GPS',
       '• **Position** — GPS to the search area, **vision** owns the final metre (no RTK needed)',
       '• **Control ×3** — CM4 mission, PX4 failsafes, human pilot',
       '• **Link loss** — mission continues on the CM4; FC returns if the datalink drops',
       '• **Release gates** — id-verified land + touchdown-gated release, both audited'],820,128,400,560,12),
],'ความปลอดภัย 3 ชั้น: นักบิน (POSCTL → companion ถอย), PX4 failsafe (geofence/datalink/RC/แบต → Return), companion watchdog (เพดาน/พื้น/no-fly/แบต/GPS). checklist ก่อนบิน + อ่าน 24 พารามิเตอร์จากบอร์ด. redundancy: 6 ใบพัด, แบต 2 ก้อน, ความสูง 3 แหล่ง, vision คุมเมตรสุดท้าย')

# 9a PROBLEMS 1-6
slide('10 · PROBLEMS WE HIT (1/2)','Starting from zero — team, budget and the first airframes',[
  table([['#','Problem','What we did'],
         ['1','No prior experience — none of us had built an autonomous drone','Four months of preparation from zero: rules → SITL → hardware → field'],
         ['2','No funding','Students and advisor paid out of pocket; asked industry for an aircraft; reused parts left by senior projects'],
         ['3','First month: hardware not ready — time must not be wasted','Built the software first: the GCS and the SITL, so the mission code was flying in simulation before any aircraft'],
         ['4','Tried Pixhawk 6X + CM4 on a small drone first','Could not even hover — frame and motors did not match'],
         ['5','Moved to a larger carbon frame','Still not flyable'],
         ['6','Needed a real airframe','DRON (Defence Research Operation Network) lent us the EFT X6100 frame + ESCs + propellers — flies well and efficiently']],60,128,760,562,[30,330,400],12),
  ph('รูปโดรนลำเล็ก (item 4)',840,128,380,175,'small test drone'),
  ph('รูปเฟรมคาร์บอน (item 5)',840,315,380,175,'carbon frame'),
  ph('รูป EFT X6100 จาก DRON (item 6)',840,502,380,188,'the aircraft we fly'),
],'ปัญหาช่วงแรก: ไม่มีประสบการณ์ (เตรียม 4 เดือนจากศูนย์), ไม่มีงบ (ออกเงินเอง/ขอเอกชน/ของรุ่นพี่), เดือนแรกฮาร์ดแวร์ไม่พร้อมเลยทำ GCS+SITL ก่อน, ลองโดรนเล็กแล้ว hover ไม่ได้, เฟรมคาร์บอนก็ไม่ได้, จนได้ EFT X6100 จาก DRON')

# 9b PROBLEMS 7-13
slide('11 · PROBLEMS WE HIT (2/2)','Hardware and field issues — found by flying, fixed by measuring',[
  table([['#','Problem','What we did'],
         ['7','7 500 mAh pack could not finish the mission','DRON added 17 000 and 15 000 mAh packs'],
         ['8','GPS freezing / unstable','Replaced the GPS cable'],
         ['9','Camera: shake, over-exposure, low frame rate','Camera settings + our grabber code: MJPEG passthrough (25 fps), highlight-priority auto-exposure that meters the pad, not the grass'],
         ['10','PM03D power module failed — wrong voltage reading','Switched to the PM02D'],
         ['11','PM02D has no supply for the CM4','CM4 runs from a power bank'],
         ['12','Holybro SiK telemetry died; RC had no long-range module','RadioMaster NOMAD — long-range RC and MAVLink telemetry in one'],
         ['13','Telemetry link unreliable','Found the soldered RX signal wire broken — re-soldered']],60,128,700,562,[30,300,370],12),
  img(A+'/real_decodes.jpg',780,128,440,248,'Item 9 fixed — real flight 26 Aug, 8 m AGL: ids 1·5·2·4 decoded live'),
  ph('รูป NOMAD / PM02D / สาย GPS (รอถ่าย)',780,395,440,140,'items 8 · 10 · 12'),
  txt(['**Lesson:** the simulator never showed any of these — flying for real, then reading every flight log, did.'],780,550,440,140,13),
],'ปัญหาฮาร์ดแวร์: แบต 7500 ไม่พอ → DRON ให้ 17000/15000, GPS ค้าง → เปลี่ยนสาย, กล้องสั่น/แสงเกิน/เฟรมเรตต่ำ → ปรับ setting + โค้ดกล้อง (ภาพขวาคือหลังแก้ decode ได้กลางอากาศ), PM03D เสีย → PM02D, PM02D ไม่จ่าย CM4 → powerbank, SiK พัง → NOMAD, สาย RX ขาด → บัดกรีใหม่')

# 10 REAL FLIGHTS
slide('12 · FLYING FOR REAL','KMUTNB practice field — five field days, a dozen-plus flights',[
  video('media/real_flight_2026-08-26_decoded.mp4',A+'/real_flight3_poster.jpg',60,128,640,360,'Real flight 26 Aug 2026 — nadir camera with live ArUco decode (green = id read)'),
  tiles([('5','field days','20–26 Aug at KMUTNB','acc'),('12+','real flights','every ULog audited','ink'),('3/3','transit points','24 Aug flight','go'),('4 ids','decoded in flight','8 m AGL · 26 Aug','amb')],720,128,500,170,2),
  txt(['• Transit corridor, geofence, pilot takeover and the camera pipeline are each **proven on the aircraft**',
       '• The full chain — sweep → land on pad → release — is complete in SITL; the **28 Aug flight trial** is the dress rehearsal',
       '• 26 Aug alone found three things no simulator shows: PX4 rewriting home altitude in flight, a mission clock fed by a sparse GPS timestamp, and auto-exposure metering the grass — **all fixed**'],720,310,500,380,13),
  ph('รูปทีมที่สนาม / โดรนขณะบิน (รอถ่าย)',60,505,640,185,'field photo'),
],'บินจริงที่สนาม มจพ. 5 วัน กว่าสิบเที่ยว วิดีโอคือกล้อง nadir ขณะบินวันที่ 26 ส.ค. ที่ decode ArUco ได้สด แต่ละส่วนพิสูจน์บนเครื่องจริงแล้ว; ลูปเต็มครบใน SITL; วันศุกร์ trial slot คือ dress rehearsal')

# 11 CLOSING
slide('AEROOPTIX · KMUTNB','Thank you — questions?',[
  tiles([('Concept','assigned-id delivery · mandatory corridor · survey once, then serve · land-on-pad release','','acc'),('Engineering','hexacopter with margin · PX4 + MAVLink · our own SITL and GCS · vision owns the final metre','','amb'),('Safety','three independent layers · pilot always wins · checklist + 24-param board read · redundancy in rotors, packs, sensors','','go')],60,150,1160,230,3),
  txt(['## Deterministic. Offline. Full-auto — inside 20 minutes.','','Team AeroOptix · King Mongkut\'s University of Technology North Bangkok · AAVC 2026, IAAI KMITL'],60,430,1160,200,18),
],'สรุป 3 ด้าน แล้วเปิด Q&A','close')

# ---------------- helpers ----------------
def b64(path):
    ext=path.rsplit('.',1)[-1].lower(); mime='image/png' if ext=='png' else 'image/jpeg'
    return f'data:{mime};base64,'+base64.b64encode(open(path,'rb').read()).decode()
def fit(path,w,h):
    iw,ih=Image.open(path).size; s=min(w/iw,h/ih); return int(iw*s),int(ih*s)
def inline(s):
    s=H.escape(s); s=re.sub(r'\*\*(.+?)\*\*',r'<b>\1</b>',s); return s

# ---------------- HTML ----------------
fonts=open(S+'/fonts_style.html',encoding='utf-8').read()
css=f"""
:root{{--bg:#{C['bg']};--surf:#{C['surf']};--line:#{C['line']};--line2:#{C['line2']};--ink:#{C['ink']};--mut:#{C['mut']};--dim:#{C['dim']};--acc:#{C['acc']};--amb:#{C['amb']};--go:#{C['go']};--crit:#{C['crit']};
--sans:'IBM Plex Sans',system-ui,sans-serif;--mono:'IBM Plex Mono',monospace;--disp:'Chakra Petch','IBM Plex Sans',sans-serif}}
*{{box-sizing:border-box;margin:0;padding:0}} html,body{{height:100%;background:var(--bg);color:var(--ink);font-family:var(--sans);overflow:hidden}}
#wrap{{position:fixed;inset:0}}
.stage{{position:absolute;left:0;top:0;width:1280px;height:720px;transform-origin:top left;background:var(--bg);overflow:hidden}}
.stage::before{{content:"";position:absolute;inset:0;pointer-events:none;background:linear-gradient(90deg,rgba(51,227,196,.03) 1px,transparent 1px) 0 0/64px 64px,linear-gradient(rgba(51,227,196,.03) 1px,transparent 1px) 0 0/64px 64px,radial-gradient(900px 500px at 80% -10%,rgba(51,227,196,.09),transparent 60%)}}
.slide{{position:absolute;inset:0;display:none}} .slide.on{{display:block}}
.eyebrow{{position:absolute;left:60px;top:38px;font-family:var(--mono);font-size:12px;letter-spacing:.28em;text-transform:uppercase;color:var(--acc)}}
.eyebrow::before{{content:"";display:inline-block;width:22px;height:2px;background:var(--acc);margin-right:10px;vertical-align:middle;box-shadow:0 0 8px var(--acc)}}
h1{{position:absolute;left:60px;top:58px;width:1160px;font-family:var(--disp);font-weight:600;font-size:34px;line-height:1.1;letter-spacing:-.005em}}
.slide.title h1{{display:none}}
.blk{{position:absolute}}
.img{{display:flex;flex-direction:column;align-items:center;justify-content:flex-start}} .img img{{display:block;border-radius:8px;border:1px solid var(--line)}}
.cap{{font-family:var(--mono);font-size:10.5px;letter-spacing:.06em;text-transform:uppercase;color:var(--dim);text-align:center;margin-top:6px;line-height:1.3}}
.ph{{border:2px dashed var(--line2);border-radius:12px;background:rgba(20,29,43,.6);display:flex;flex-direction:column;align-items:center;justify-content:center;gap:6px;color:var(--mut);text-align:center;padding:10px}}
.ph .i{{font-size:34px;opacity:.6}} .ph .t{{font-size:15px;font-weight:600;color:var(--ink)}} .ph .s{{font-family:var(--mono);font-size:10.5px;letter-spacing:.08em;text-transform:uppercase;color:var(--dim)}}
.txt{{color:var(--mut);line-height:1.45}} .txt p{{margin:0 0 6px}} .txt b{{color:var(--ink);font-weight:600}} .txt .h{{font-family:var(--disp);font-weight:600;color:var(--acc);font-size:1.12em;margin:8px 0 4px}} .txt .h:first-child{{margin-top:0}}
.txt .li{{padding-left:18px;position:relative}} .txt .li::before{{content:"▸";position:absolute;left:0;color:var(--acc);font-size:.85em;top:1px}}
.tiles{{display:grid;gap:12px}} .tile{{background:var(--surf);border:1px solid var(--line);border-radius:10px;padding:12px 14px;display:flex;flex-direction:column;justify-content:center;min-height:0}}
.tile .v{{font-family:var(--mono);font-weight:600;font-size:26px;line-height:1.05}} .tile .k{{font-family:var(--mono);font-size:10.5px;letter-spacing:.12em;text-transform:uppercase;color:var(--dim);margin-top:6px}} .tile .s{{font-size:12px;color:var(--mut);margin-top:3px}}
.tile.big .v{{font-family:var(--disp);font-size:22px}} .tile.big .k{{font-family:var(--sans);text-transform:none;letter-spacing:0;font-size:13.5px;color:var(--mut);line-height:1.4;margin-top:8px}}
.c-acc{{color:var(--acc)}} .c-amb{{color:var(--amb)}} .c-go{{color:var(--go)}} .c-ink{{color:var(--ink)}}
table{{border-collapse:collapse;width:100%;color:var(--mut)}} th{{font-family:var(--mono);font-size:10.5px;letter-spacing:.12em;text-transform:uppercase;color:var(--acc);text-align:left;padding:6px 8px;border-bottom:1px solid var(--line2)}}
td{{padding:6px 8px;border-bottom:1px solid var(--line);vertical-align:top;line-height:1.35}} td:first-child{{font-family:var(--mono);color:var(--amb)}} tr td:nth-child(2){{color:var(--ink)}}
.video{{display:flex;flex-direction:column;align-items:center}} video{{display:block;border-radius:8px;border:1px solid var(--line);background:#000}}
.hud{{position:fixed;z-index:9;font-family:var(--mono);font-size:11px;letter-spacing:.16em;color:var(--dim);text-transform:uppercase}} .hud.tl{{left:22px;top:16px}} .hud.tr{{right:22px;top:16px}} .hud.br{{right:22px;bottom:14px}}
.prog{{position:fixed;top:0;left:0;height:3px;background:var(--acc);z-index:10;box-shadow:0 0 10px var(--acc);transition:width .3s}}
#notes{{position:fixed;left:0;right:0;bottom:0;background:rgba(10,14,21,.96);border-top:1px solid var(--line2);color:var(--mut);padding:14px 26px;font-size:15px;line-height:1.5;display:none;z-index:20}} #notes.on{{display:block}}
.title .big{{font-family:var(--disp);font-weight:700;font-size:74px;color:var(--ink);line-height:1}}
"""
parts=[]
for n,sl in enumerate(SL,1):
    inner=[f'<div class="eyebrow">{H.escape(sl["eyebrow"])}</div>']
    if sl['kind']=='title':
        inner.append('<div class="blk" style="left:60px;top:150px;width:640px"></div>')
    else:
        inner.append(f'<h1>{H.escape(sl["title"])}</h1>')
    if sl['kind']=='close':
        inner.append(f'<div class="blk" style="left:60px;top:70px;width:1160px"><div class="title big" style="font-size:56px">{H.escape(sl["title"])}</div></div>')
        inner=[x for x in inner if not x.startswith('<h1')]
    for b in sl['blocks']:
        st=f'left:{b["x"]}px;top:{b["y"]}px;width:{b["w"]}px;height:{b["h"]}px'
        if b['t']=='img':
            ch=b['h']-(24 if b['cap'] else 0); w,h=fit(b['src'],b['w'],ch)
            inner.append(f'<div class="blk img" style="{st}"><img src="{b64(b["src"])}" style="width:{w}px;height:{h}px" alt="">'+(f'<div class="cap">{H.escape(b["cap"])}</div>' if b['cap'] else '')+'</div>')
        elif b['t']=='ph':
            inner.append(f'<div class="blk ph" style="{st}"><div class="i">📷</div><div class="t">{H.escape(b["text"])}</div>'+(f'<div class="s">{H.escape(b["sub"])}</div>' if b['sub'] else '')+'</div>')
        elif b['t']=='txt':
            ps=[]
            for ln in b['lines']:
                if ln.startswith('## '): ps.append(f'<p class="h">{inline(ln[3:])}</p>')
                elif ln.startswith('• '): ps.append(f'<p class="li">{inline(ln[2:])}</p>')
                elif ln=='' : ps.append('<p style="height:6px"></p>')
                else: ps.append(f'<p>{inline(ln)}</p>')
            if sl['kind']=='title' and b is sl['blocks'][0]:
                ps[0]=f'<div class="big">AeroOptix</div>'
            inner.append(f'<div class="blk txt" style="{st};font-size:{b["size"]}px">'+''.join(ps)+'</div>')
        elif b['t']=='tiles':
            big=sl['kind']=='close'
            cells=''.join(f'<div class="tile{" big" if big else ""}"><div class="v c-{c}">{H.escape(v)}</div><div class="k">{H.escape(k)}</div>'+(f'<div class="s">{H.escape(s2)}</div>' if s2 else '')+'</div>' for v,k,s2,c in b['items'])
            inner.append(f'<div class="blk tiles" style="{st};grid-template-columns:repeat({b["cols"]},1fr)">{cells}</div>')
        elif b['t']=='video':
            ch=b['h']-(24 if b['cap'] else 0); w=b['w']; h=min(ch,int(w*9/16)); w=int(h*16/9)
            inner.append(f'<div class="blk video" style="{st}"><video controls preload="metadata" poster="{b64(b["poster"])}" src="{b["src"]}" style="width:{w}px;height:{h}px"></video>'+(f'<div class="cap">{H.escape(b["cap"])}</div>' if b['cap'] else '')+'</div>')
        elif b['t']=='table':
            hdr,*rows=b['rows']; cols=''.join(f'<col style="width:{cw}px">' for cw in b['colw'])
            th=''.join(f'<th>{H.escape(c)}</th>' for c in hdr)
            tr=''.join('<tr>'+''.join(f'<td>{inline(c)}</td>' for c in r)+'</tr>' for r in rows)
            inner.append(f'<div class="blk" style="{st};font-size:{b["size"]}px;overflow:hidden"><table><colgroup>{cols}</colgroup><thead><tr>{th}</tr></thead><tbody>{tr}</tbody></table></div>')
    notes=H.escape(sl['notes'])
    parts.append(f'<section class="slide {sl["kind"]}" data-notes="{notes}">'+''.join(inner)+'</section>')
js="""
(function(){var sl=[].slice.call(document.querySelectorAll('.slide')),N=sl.length,i=0,st=document.getElementById('stage'),prog=document.getElementById('prog'),hud=document.getElementById('hudno'),notes=document.getElementById('notes');
function fitS(){var s=Math.min(window.innerWidth/1280,window.innerHeight/720);st.style.transform='scale('+s+')';st.style.marginLeft=((window.innerWidth-1280*s)/2)+'px';st.style.marginTop=((window.innerHeight-720*s)/2)+'px';}
function show(k){k=Math.max(0,Math.min(N-1,k));sl.forEach(function(x,j){x.classList.toggle('on',j==k);if(j!=k)x.querySelectorAll('video').forEach(function(v){v.pause()})});i=k;prog.style.width=(k/(N-1)*100)+'%';hud.textContent='SLIDE '+(k+1<10?'0':'')+(k+1)+' / '+N;notes.textContent=sl[k].dataset.notes;location.hash='s'+(k+1);}
window.addEventListener('resize',fitS);window.addEventListener('keydown',function(e){if(e.target.tagName==='VIDEO')return;
if(['ArrowRight','ArrowDown',' ','PageDown'].indexOf(e.key)>=0){e.preventDefault();show(i+1)}else if(['ArrowLeft','ArrowUp','PageUp'].indexOf(e.key)>=0){e.preventDefault();show(i-1)}else if(e.key==='Home'){show(0)}else if(e.key==='End'){show(N-1)}else if(e.key==='n'||e.key==='N'){notes.classList.toggle('on')}else if(e.key==='f'||e.key==='F'){document.documentElement.requestFullscreen&&document.documentElement.requestFullscreen()}});
document.getElementById('wrap').addEventListener('click',function(e){if(e.target.closest('video'))return;show(e.clientX>window.innerWidth/2?i+1:i-1)});
var m=/s(\\d+)/.exec(location.hash);fitS();show(m?parseInt(m[1])-1:0);})();
"""
page=f"""<!doctype html><html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>AeroOptix — AAVC 2026</title>{fonts}<style>{css}</style></head><body>
<div class="prog" id="prog"></div><div class="hud tl">AAVC 2026 · <span style="color:var(--acc)">TEAM AEROOPTIX</span></div><div class="hud tr" id="hudno"></div><div class="hud br">→ / SPACE NEXT · ← BACK · N NOTES · F FULLSCREEN</div>
<div id="wrap"><div class="stage" id="stage">{''.join(parts)}</div></div><div id="notes"></div><script>{js}</script></body></html>"""
open(OUT+'/AAVC2026_Presentation_v2.html','w',encoding='utf-8').write(page)
print('html bytes',len(page))

# ---------------- PPTX ----------------
prs=Presentation(); prs.slide_width=Emu(1280*9525); prs.slide_height=Emu(720*9525)
blank=prs.slide_layouts[6]
def E(px): return Emu(int(px*9525))
def rect(sl,x,y,w,h,fill=None,line=None,shape=MSO_SHAPE.RECTANGLE,dash=False,lw=1):
    sh=sl.shapes.add_shape(shape,E(x),E(y),E(w),E(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill)
    else: sh.fill.background()
    if line: sh.line.color.rgb=rgb(line); sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    if dash: sh.line.dash_style=MSO_LINE.DASH
    sh.shadow.inherit=False
    return sh
def runs_from(par,text,size,color='mut',bold_color='ink',font='Calibri',bold_all=False):
    for seg in re.split(r'(\*\*.+?\*\*)',text):
        if not seg: continue
        r=par.add_run(); b=seg.startswith('**'); r.text=seg.strip('*') if b else seg
        r.font.size=Pt(size); r.font.name=font; r.font.bold=b or bold_all; r.font.color.rgb=rgb(bold_color if b else color)
def tbox(sl,x,y,w,h,lines,size,color='mut',anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT,font='Calibri'):
    tb=sl.shapes.add_textbox(E(x),E(y),E(w),E(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Emu(0); tf.margin_top=tf.margin_bottom=Emu(0)
    first=True
    for ln in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        if isinstance(ln,tuple):  # (text,size,color,bold,space_after)
            t,sz,col,bd,sa=ln; runs_from(p,t,sz,col,col if bd else 'ink',font,bold_all=bd); p.space_after=Pt(sa)
        elif ln.startswith('## '): runs_from(p,ln[3:],size*1.12,'acc','acc',font,bold_all=True); p.space_after=Pt(2); p.space_before=Pt(4)
        elif ln.startswith('• '): runs_from(p,'▸ '+ln[2:],size,color,'ink',font); p.space_after=Pt(3)
        elif ln=='': p.space_after=Pt(4)
        else: runs_from(p,ln,size,color,'ink',font); p.space_after=Pt(3)
    return tb
def pic(sl,path,x,y,w,h,cap=None):
    ch=h-(22 if cap else 0); pw,ph_=fit(path,w,ch); px=x+(w-pw)//2
    sl.shapes.add_picture(path,E(px),E(y),E(pw),E(ph_))
    if cap: tbox(sl,x,y+ph_+4,w,20,[(cap,7.5,'dim',False,0)],7.5,'dim',align=PP_ALIGN.CENTER,font='Consolas')
for n,sl in enumerate(SL,1):
    s=prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb('bg')
    tbox(s,60,34,900,18,[(sl['eyebrow'],9,'acc',False,0)],9,'acc',font='Consolas')
    if sl['kind']=='title':
        tbox(s,60,120,700,80,[('AeroOptix',54,'ink',True,0)],54,'ink')
    elif sl['kind']=='close':
        tbox(s,60,66,1160,60,[(sl['title'],40,'ink',True,0)],40,'ink')
    else:
        tbox(s,60,58,1160,60,[(sl['title'],24,'ink',True,0)],24,'ink')
    for b in sl['blocks']:
        if b['t']=='img': pic(s,b['src'],b['x'],b['y'],b['w'],b['h'],b['cap'])
        elif b['t']=='ph':
            sh=rect(s,b['x'],b['y'],b['w'],b['h'],'bg2','line2',MSO_SHAPE.ROUNDED_RECTANGLE,dash=True,lw=1.5)
            tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; runs_from(p,'📷  '+b['text'],13,'ink','ink',font='Tahoma',bold_all=True)
            if b['sub']: p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; runs_from(p2,b['sub'],8,'dim','dim',font='Consolas')
        elif b['t']=='txt':
            lines=b['lines']
            if sl['kind']=='title' and b is sl['blocks'][0]: lines=lines[1:]; y=b['y']+60
            else: y=b['y']
            tbox(s,b['x'],y,b['w'],b['h'],lines,b['size']*0.72,font='Tahoma' if any(ord(ch)>0x0E00 and ord(ch)<0x0E7F for ln in lines for ch in ln) else 'Calibri')
        elif b['t']=='tiles':
            cols=b['cols']; rows=(len(b['items'])+cols-1)//cols; gap=12
            tw=(b['w']-gap*(cols-1))/cols; th=(b['h']-gap*(rows-1))/rows
            for k,(v,kk,s2,c) in enumerate(b['items']):
                cx=b['x']+(k%cols)*(tw+gap); cy=b['y']+(k//cols)*(th+gap)
                sh=rect(s,cx,cy,tw,th,'surf','line',MSO_SHAPE.ROUNDED_RECTANGLE); sh.adjustments[0]=0.08
                tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=tf.margin_right=E(12)
                big=sl['kind']=='close'
                p=tf.paragraphs[0]; runs_from(p,v,16 if big else 19,c,c,font='Calibri' if big else 'Consolas',bold_all=True); p.space_after=Pt(3)
                p=tf.add_paragraph(); runs_from(p,kk,10 if big else 7.5,'mut' if big else 'dim','ink',font='Calibri' if big else 'Consolas')
                if s2: p=tf.add_paragraph(); runs_from(p,s2,8.5,'mut','ink')
        elif b['t']=='video':
            ch=b['h']-(24 if b['cap'] else 0); w=b['w']; h=min(ch,int(w*9/16)); w=int(h*16/9)
            src=OUT+'/'+b['src']
            s.shapes.add_movie(src,E(b['x']),E(b['y']),E(w),E(h),poster_frame_image=b['poster'],mime_type='video/mp4')
            if b['cap']: tbox(s,b['x'],b['y']+h+4,b['w'],20,[(b['cap'],7.5,'dim',False,0)],7.5,'dim',align=PP_ALIGN.CENTER,font='Consolas')
        elif b['t']=='table':
            rows=b['rows']; nr=len(rows); nc=len(rows[0])
            tb=s.shapes.add_table(nr,nc,E(b['x']),E(b['y']),E(b['w']),E(b['h'])).table
            for j,cw in enumerate(b['colw']): tb.columns[j].width=E(cw*b['w']/sum(b['colw']))
            for i,row in enumerate(rows):
                tb.rows[i].height=E(b['h']/nr)
                for j,cell in enumerate(row):
                    c=tb.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=rgb('surf' if i%2 else 'bg2'); c.margin_left=c.margin_right=E(6); c.margin_top=c.margin_bottom=E(3)
                    tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
                    if i==0: runs_from(p,cell,7.5,'acc','acc',font='Consolas',bold_all=True)
                    elif j==0: runs_from(p,cell,b['size']*0.9,'amb','amb',font='Consolas')
                    elif j==1: runs_from(p,cell,b['size']*0.9,'ink','ink')
                    else: runs_from(p,cell,b['size']*0.9,'mut','ink')
    ntf=s.notes_slide.notes_text_frame; ntf.text=sl['notes']
    for pp in ntf.paragraphs:
        for r in pp.runs: r.font.name='Tahoma'; r.font.size=Pt(12)
prs.save(OUT+'/AAVC2026_Presentation_v2.pptx'); print('pptx ok', os.path.getsize(OUT+'/AAVC2026_Presentation_v2.pptx'))
